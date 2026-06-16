"""
Generate one PPTX training presentation per role for the KSB Purchase
Requisition System.  Run with:  python generate_pptx.py
Output files land next to this script.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── Colours ──────────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x06, 0x2e, 0x5c)   # deep navy
C_BLUE   = RGBColor(0x0b, 0x4a, 0x8a)   # KSB primary
C_MID    = RGBColor(0x1a, 0x6f, 0xc7)   # lighter blue
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_NEAR   = RGBColor(0x11, 0x11, 0x11)
C_MUTED  = RGBColor(0x55, 0x55, 0x55)
C_SOFT   = RGBColor(0xea, 0xf2, 0xfb)   # very light blue tint
C_BADGE  = RGBColor(0xf0, 0xb4, 0x29)   # amber accent for badges
C_GREEN  = RGBColor(0x06, 0x5f, 0x46)

# ── Dimensions ───────────────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
LM = Inches(0.55)   # left margin
RM = Inches(0.55)   # right margin
CW = W - LM - RM    # content width

HEADER_H  = Inches(1.15)
CONTENT_Y = Inches(1.3)
CONTENT_H = H - CONTENT_Y - Inches(0.25)

# ── Low-level helpers ─────────────────────────────────────────────────────────
def rgb_hex(c): return f"{c[0]:02x}{c[1]:02x}{c[2]:02x}"

def _set_solid(shape, color):
    """Fill a shape with a solid colour."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    _set_solid(shape, color)
    return shape

def _tf(shape, text, size, bold=False, color=C_NEAR, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def _add_textbox(slide, x, y, w, h, text, size=14, bold=False,
                 color=C_NEAR, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    _tf(tb, text, size, bold=bold, color=color, align=align, italic=italic)
    return tb

def _para(tf, text, size=13, bold=False, color=C_NEAR,
          align=PP_ALIGN.LEFT, space_before=Pt(2), bullet=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

# ── Slide builders ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def cover_slide(prs, role_label, subtitle, meta_rows):
    """Full-bleed navy cover with white text and inverted role badge."""
    sl = blank_slide(prs)

    # Background gradient: two overlapping rects (dark → mid)
    _add_rect(sl, 0, 0, W, H, C_DARK)
    _add_rect(sl, 0, 0, W, H*0.55, C_BLUE)

    # Top stripe
    _add_rect(sl, 0, 0, W, Inches(0.08), RGBColor(0xff, 0xff, 0xff))
    # Bottom stripe
    _add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), RGBColor(0xff, 0xff, 0xff))

    # Kicker
    tb = slide.shapes.add_textbox if False else sl.shapes.add_textbox
    k = sl.shapes.add_textbox(LM, Inches(1.0), CW, Inches(0.4))
    k.text_frame.word_wrap = False
    _tf(k, "KSB ZAMBIA  ·  PURCHASE REQUISITION SYSTEM",
        9, color=RGBColor(0xb0, 0xc8, 0xe8), align=PP_ALIGN.CENTER)

    # Main title
    t = sl.shapes.add_textbox(LM, Inches(1.6), CW, Inches(1.0))
    _tf(t, "User Manual", 38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Role badge (white rounded rect with blue text)
    bw, bh = Inches(5.5), Inches(0.75)
    bx = (W - bw) / 2
    by = Inches(2.85)
    badge = sl.shapes.add_shape(1, bx, by, bw, bh)
    badge.line.fill.background()
    _set_solid(badge, C_WHITE)
    btf = badge.text_frame
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = role_label
    br.font.size = Pt(22)
    br.font.bold = True
    br.font.color.rgb = C_BLUE

    # Subtitle
    sub = sl.shapes.add_textbox(LM, Inches(3.85), CW, Inches(0.55))
    _tf(sub, subtitle, 14, color=RGBColor(0xb0, 0xc8, 0xe8), align=PP_ALIGN.CENTER)

    # Meta table (plain text rows)
    my = Inches(4.8)
    for label, value in meta_rows:
        row = sl.shapes.add_textbox(Inches(3.2), my, Inches(7.0), Inches(0.32))
        tf = row.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = f"{label}:  "; r1.font.size = Pt(11)
        r1.font.bold = True; r1.font.color.rgb = RGBColor(0x8a, 0xb4, 0xd8)
        r2 = p.add_run(); r2.text = value; r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(0xd0, 0xe3, 0xf2)
        my += Inches(0.33)

    return sl


def content_slide(prs, title, bullets, two_col=False):
    """Standard slide: blue header bar + bulleted content."""
    sl = blank_slide(prs)

    # Header bar
    _add_rect(sl, 0, 0, W, HEADER_H, C_BLUE)
    th = sl.shapes.add_textbox(LM, Inches(0.3), CW, Inches(0.75))
    _tf(th, title, 22, bold=True, color=C_WHITE)

    return sl


def _bullet_block(slide, bullets, x, y, w, h, title=None):
    """Add a block of bullet points, optionally with a bold heading."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if item.startswith("##"):          # section heading
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.space_before = Pt(6)
            r = p.add_run()
            r.text = item[2:].strip()
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = C_BLUE
        elif item.startswith("--"):        # sub-bullet
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.level = 1
            p.space_before = Pt(1)
            r = p.add_run()
            r.text = item[2:].strip()
            r.font.size = Pt(12)
            r.font.color.rgb = C_MUTED
        else:
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            p.level = 0
            p.space_before = Pt(3)
            r = p.add_run()
            r.text = item
            r.font.size = Pt(13)
            r.font.color.rgb = C_NEAR
    return tb


def bullets_slide(prs, title, items, note=None):
    sl = content_slide(prs, title, items)
    _bullet_block(sl, items, LM, CONTENT_Y, CW, CONTENT_H)
    if note:
        nb = sl.shapes.add_textbox(LM, H - Inches(0.55), CW, Inches(0.4))
        _tf(nb, f"NOTE: {note}", 10, italic=True, color=C_MUTED)
    return sl


def two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    sl = content_slide(prs, title, [])
    col_w = (CW - Inches(0.3)) / 2

    # Left column header
    lh = sl.shapes.add_textbox(LM, CONTENT_Y, col_w, Inches(0.38))
    _tf(lh, left_title, 13, bold=True, color=C_BLUE)
    _bullet_block(sl, left_items, LM, CONTENT_Y + Inches(0.4), col_w,
                  CONTENT_H - Inches(0.4))

    # Divider
    _add_rect(sl, LM + col_w + Inches(0.12), CONTENT_Y,
              Inches(0.03), CONTENT_H - Inches(0.1),
              RGBColor(0xcc, 0xdd, 0xee))

    # Right column header
    rx = LM + col_w + Inches(0.3)
    rh = sl.shapes.add_textbox(rx, CONTENT_Y, col_w, Inches(0.38))
    _tf(rh, right_title, 13, bold=True, color=C_BLUE)
    _bullet_block(sl, right_items, rx, CONTENT_Y + Inches(0.4), col_w,
                  CONTENT_H - Inches(0.4))
    return sl


def steps_slide(prs, title, steps, intro=None):
    """Numbered step blocks."""
    sl = content_slide(prs, title, [])
    y = CONTENT_Y
    if intro:
        ib = sl.shapes.add_textbox(LM, y, CW, Inches(0.35))
        _tf(ib, intro, 12, italic=True, color=C_MUTED)
        y += Inches(0.38)

    step_h = min(Inches(0.62), (H - y - Inches(0.2)) / max(len(steps), 1))
    for i, step in enumerate(steps, 1):
        # Circle number
        circ = sl.shapes.add_shape(9, LM, y + Inches(0.06),
                                   Inches(0.45), Inches(0.45))   # oval
        _set_solid(circ, C_BLUE)
        circ.line.fill.background()
        ctf = circ.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i)
        cr.font.size = Pt(12)
        cr.font.bold = True
        cr.font.color.rgb = C_WHITE

        # Step text
        stb = sl.shapes.add_textbox(LM + Inches(0.58), y,
                                    CW - Inches(0.6), step_h)
        stb.text_frame.word_wrap = True
        _tf(stb, step, 13, color=C_NEAR)
        y += step_h
    return sl


def workflow_slide(prs, title, lines):
    """Monospace workflow diagram."""
    sl = content_slide(prs, title, [])
    # Light background box
    _add_rect(sl, LM, CONTENT_Y, CW, CONTENT_H - Inches(0.1), C_SOFT)
    tb = sl.shapes.add_textbox(LM + Inches(0.3), CONTENT_Y + Inches(0.2),
                               CW - Inches(0.6), CONTENT_H - Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.name = "Consolas"
        r.font.color.rgb = C_NEAR
    return sl


def end_slide(prs, role_label):
    sl = blank_slide(prs)
    _add_rect(sl, 0, 0, W, H, C_DARK)
    _add_rect(sl, 0, 0, W, Inches(0.08), C_WHITE)
    _add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), C_WHITE)

    tb = sl.shapes.add_textbox(LM, Inches(2.4), CW, Inches(1.0))
    _tf(tb, "Refer to your role manual for full detail.", 20,
        color=C_WHITE, align=PP_ALIGN.CENTER)

    tb2 = sl.shapes.add_textbox(LM, Inches(3.6), CW, Inches(0.6))
    _tf(tb2, f"{role_label}  ·  KSB Zambia Purchase Requisition System",
        13, color=RGBColor(0x8a, 0xb4, 0xd8), align=PP_ALIGN.CENTER)
    return sl


# ── Role definitions ──────────────────────────────────────────────────────────

ROLES = [

    {
        "file": "00-all-users",
        "label": "All Users — Orientation",
        "subtitle": "Login · Navigation · Security basics",
        "meta": [("Audience", "Everyone — read before your role manual"),
                 ("Version", "1.0")],
        "slides": [
            ("bullets", "What the System Does", [
                "## Document types you can work with:",
                "Purchase Requisitions — request goods or services",
                "EFT Requisitions — request a bank transfer payment",
                "Petty Cash Requisitions — draw small cash amounts",
                "Expense Claims — get reimbursed for money you spent",
                "## Stores module (warehouse staff only):",
                "Goods Receipt Notes (GRNs) — record incoming stock",
                "Issue Slips — record stock leaving the warehouse",
                "Picking Slips — record stock moving between locations",
                "Real-Time Stock Register — live view of available stock",
            ]),
            ("workflow", "Approval Flow at a Glance", [
                "Purchase Requisition:",
                "",
                "  [Initiator] --> pending_hod --> pending_finance --> pending_md* --> approved",
                "",
                "  * pending_md only when value exceeds the MD threshold.",
                "",
                "EFT / Petty Cash / Expense Claim / Issue Slip:",
                "",
                "  [Initiator] --> pending_hod --> pending_finance --> approved",
                "",
                "GRN:    [Stores] --> pending_approval --> approved    (Finance only)",
                "Picking Slip:    No approval — created as Completed immediately.",
            ]),
            ("steps", "First Login — Forced Password Change", [
                "Enter your username and temporary password on the login screen.",
                "The system intercepts login and shows the Change Password screen — you cannot skip it.",
                "Enter your current temporary password.",
                "Enter and confirm your new password (minimum 8 characters, must differ from the temporary one).",
                "Select and answer 3 security questions — used for future self-service resets.",
                "Click Save. You are taken to the Dashboard automatically.",
            ], "This only happens on your very first login or after an admin resets your password."),
            ("two_col", "Screen Layout", [
                "## Sidebar (left)",
                "Main menu, grouped by section",
                "Items hidden if your role has no access",
                "Click a group header to expand / collapse",
                "## Top Bar",
                "Your name and role on the right",
                "Click your name → Profile / Logout",
                "Breadcrumb shows where you are",
            ], [
                "## Main Area (right)",
                "Shows the page you clicked: list, form, or detail",
                "Download PDF / Print buttons at the top of detail pages",
                "Approve / Reject at the bottom of approval pages",
                "## Session",
                "Idle timeout: ~15 minutes",
                "Always log out on shared computers",
            ]),
            ("bullets", "Security Rules", [
                "Password minimum: 8 characters, different from current password",
                "Brute-force protection: too many failed logins from one IP → 15-minute block",
                "Audit logging: every approval, rejection, and status change recorded with your name, IP, browser, and timestamp — you are accountable",
                "All data travels over HTTPS",
                "## Never share your password.",
            ]),
        ],
    },

    {
        "file": "01-initiator",
        "label": "Role 1 — Initiator",
        "subtitle": "Submit requisitions, EFTs, petty cash, and expense claims",
        "meta": [("Approval role", "None — you create, not approve"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "Create Purchase Requisitions",
                "Create EFT Requisitions",
                "Create Petty Cash requests",
                "Create Expense Claims",
                "View status of your own requests",
                "Download PDF of any submission",
            ], [
                "## You CANNOT:",
                "Approve or reject any request",
                "See other people's requests",
                "Edit a submitted request",
                "-- If wrong: ask an approver to reject, then resubmit as new",
            ]),
            ("workflow", "Approval Workflows", [
                "Purchase Requisition:",
                "  [You] --> pending_hod --> pending_finance --> pending_md* --> approved",
                "  * Only if value exceeds MD threshold.",
                "",
                "EFT / Petty Cash / Expense Claim:",
                "  [You] --> pending_hod --> pending_finance --> approved",
                "  (MD is NOT in this chain.)",
                "",
                "On rejection at any stage:",
                "  --> rejected  (you are notified by email with the reason)",
                "  You must create a new submission — you cannot edit the rejected one.",
            ]),
            ("steps", "Task: Create a Purchase Requisition", [
                "Sidebar → Procurement → Create Requisition.",
                "Fill the header: Title, Description, Department, Delivery Location, Required Date, Currency.",
                "Click Add Item for each item: Name, Quantity, Unit, Estimated Unit Price.",
                "Fill the Justification box — be specific. Vague requests are frequently rejected.",
                "Attach a supplier quote (PDF / image, max 10 MB) if you have one.",
                "Review everything, then click Submit.",
                "Status becomes pending_hod. Your HOD receives an email.",
            ], "ID format: KSB-PR-<timestamp>"),
            ("bullets", "Other Form Types — Quick Start", [
                "## EFT Requisition (pay a supplier by bank transfer)",
                "Sidebar → Financial Forms → EFT Requisitions → Create New",
                "Fill payee name, bank, account number, amount, purpose, attach invoice",
                "ID: KSB-EFT-…",
                "## Petty Cash (draw small cash)",
                "Sidebar → Financial Forms → Petty Cash → Create New",
                "Fill amount, purpose, itemised breakdown, required date",
                "ID: KSB-PC-…",
                "## Expense Claim (reimbursement)",
                "Sidebar → Financial Forms → Expense Claims → Create New",
                "Add one line per expense: date, category, description, amount, receipt photo",
                "ID: KSB-EXP-…",
            ]),
            ("steps", "Task: Track Requests & Handle Rejections", [
                "Open the relevant list: Procurement → My Requisitions, or Financial Forms → the form type.",
                "Read the Status column: pending_hod / pending_finance / pending_md / approved / rejected.",
                "Click any row to open the full detail page and Approval History.",
                "If rejected: read the rejection comment in Approval History.",
                "Create a new corrected submission. In the description, reference the rejected ID.",
                "Do NOT resubmit without fixing the issue — repeated unchanged submissions damage trust.",
            ]),
        ],
    },

    {
        "file": "02-hod",
        "label": "Role 2 — Head of Department",
        "subtitle": "First-stage approver for your department",
        "meta": [("Approval stage", "Stage 1 of 3 (PRs) / Stage 1 of 2 (EFT, PC, EC, Issue Slips)"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "See every request from your department",
                "Approve / reject PRs, EFTs, Petty Cash, Expense Claims, Issue Slips at Stage 1",
                "Initiate your own requests",
                "View department spend summaries",
            ], [
                "## You CANNOT:",
                "Approve requests from other departments you don't head",
                "Approve your own initiated requests",
                "Skip the Finance stage — your approval passes to Finance, not around it",
                "Modify amounts or items — reject with a comment instead",
            ]),
            ("workflow", "Workflow — Your Position", [
                "  [Initiator submits]",
                "          |",
                "          v",
                "  +---------------------------+",
                "  |  STATUS = pending_hod     |  <-- your queue",
                "  |  Your name = the approver |",
                "  +---------------------------+",
                "          |",
                "   approve / reject",
                "          |",
                "   +------+------+",
                "   |             |",
                "   v             v",
                " pending_finance  rejected",
                "                 (initiator sees comment + can resubmit)",
            ]),
            ("steps", "Task: Find & Review a Request", [
                "Dashboard: check the Pending Approvals card for your count.",
                "Sidebar → Procurement → Pending Approvals (for PRs).",
                "Sidebar → Financial Forms → Pending Approvals (for EFT / PC / EC).",
                "Sidebar → Stores → Issue Slips, filter by pending_hod (for Issue Slips).",
                "Click a row to open the detail page.",
                "Check: initiator is in your dept, items and quantities are reasonable, justification answers 'why now', quote is attached for high-value items.",
            ]),
            ("two_col", "Approve vs Reject", [
                "## To Approve:",
                "Scroll to the Comments box",
                "Write a short meaningful comment",
                "Click the green Approve button",
                "Status moves to pending_finance",
                "Finance receives an email",
                "## Target turnaround:",
                "PRs / Expense Claims: 48 hours",
                "EFT: 24 hours (time-sensitive)",
                "Petty Cash: same day",
                "Issue Slips: 24 hours",
            ], [
                "## To Reject:",
                "Write a specific rejection reason",
                "Examples:",
                "-- 'Quote required for purchases above K5,000.'",
                "-- 'Wrong department — should come from IT.'",
                "-- 'Quantity 50 seems excessive for a 10-person team.'",
                "Click the red Reject button",
                "Initiator is notified with your comment",
                "## Avoid vague reasons like 'unclear' — the initiator cannot fix those.",
            ]),
            ("bullets", "Your Pre-Approval Checklist", [
                "Request belongs to someone in my department or a direct subordinate",
                "Items and quantities make sense for the work",
                "Prices are reasonable — compare to known costs",
                "Justification clearly answers: What? Why now? Why this supplier?",
                "Quote is attached for higher-value items",
                "Department is correctly set",
                "I have left a clear, specific comment",
                "## Going on leave? Tell Admin before you go — they can reassign your queue.",
            ]),
        ],
    },

    {
        "file": "03-finance",
        "label": "Role 3 — Finance",
        "subtitle": "Second-stage approver · Budgets · FX rates · GRN approvals",
        "meta": [("Approval stage", "Stage 2 — after HOD, before MD"),
                 ("Extra", "Budgets, FX rates, GRN approvals"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "Approve / reject at pending_finance",
                "Approve GRNs (goods receipts into stock)",
                "Approve Issue Slips at Finance stage",
                "See and edit department budgets",
                "See and edit currency FX rates",
                "See every request system-wide",
                "Run financial analytics reports",
            ], [
                "## You CANNOT:",
                "Skip the HOD stage — requests must arrive via HOD approval",
                "Approve your own initiated requests",
                "Delete users or change system roles (Admin only)",
                "Modify an approved request",
                "## Your responsibility:",
                "Budget compliance, document verification, FX accuracy, prompt GRN approvals",
            ]),
            ("workflow", "Workflow — Your Position", [
                "            [HOD approved]",
                "                  |",
                "                  v",
                "      +--------------------------+",
                "      | STATUS = pending_finance |  <-- your queue",
                "      +--------------------------+",
                "                  |",
                "         approve / reject",
                "                  |",
                "       +----------+----------+",
                "       |                     |",
                "       v                     v",
                "   pending_md*           approved",
                "  (high-value PR)   (EFT, PC, EC, Issue Slip,",
                "                    or low-value PR)",
                "",
                "GRN:  [Stores creates GRN] --> pending_approval --> approved  (you approve)",
            ]),
            ("steps", "Task: Review & Approve a Financial Request", [
                "Open Pending Approvals: Procurement → Pending Approvals or Financial Forms → Pending Approvals.",
                "Click a request. Verify HOD has approved and read their comment.",
                "Check budget: Financial Planning → Budgets — does this push the dept over budget?",
                "Check supporting documents: invoice matches amount and payee details, quotes match PR prices.",
                "Check FX: if amount is in USD/EUR/GBP/ZAR, is the current ZMW rate accurate? Update if stale.",
                "Write a meaningful comment, then click Approve or Reject.",
            ]),
            ("steps", "Task: Approve a GRN", [
                "Sidebar → Stores → Goods Receipt Notes.",
                "Filter by Status: pending_approval.",
                "Open a GRN. Note the Reference PR it relates to.",
                "Compare Quantity Received against the PR's ordered quantities.",
                "Verify the supplier/vendor matches the one on the PR.",
                "Write a comment (e.g. 'Quantities verified against DN-2026-0154'), then click Approve.",
                "Stock becomes officially available in the Stock Register immediately.",
            ]),
            ("two_col", "Budgets & FX Rates", [
                "## Budgets",
                "Sidebar → Financial Planning → Budgets",
                "Lists every dept with budget and spend-to-date",
                "Click a dept → enter annual budget in ZMW → Save",
                "Year-end: reset Spent to zero, set new year budget",
                "NOTE: budget is a soft control — system does not auto-block overspend. You enforce it during approvals.",
            ], [
                "## FX Rates",
                "Sidebar → Financial Planning → FX Rates",
                "Supported currencies: USD, EUR, GBP, ZAR",
                "Click currency → enter ZMW rate → Save",
                "Recommended cadence: update every Monday using Bank of Zambia mid-rate",
                "Update ad-hoc if a rate moves by more than 2% mid-week",
                "Document big movements in a comment for MD visibility",
            ]),
        ],
    },

    {
        "file": "04-md",
        "label": "Role 4 — Managing Director",
        "subtitle": "Final approver on high-value Purchase Requisitions",
        "meta": [("Approval stage", "Stage 3 — final, for PRs above the MD threshold only"),
                 ("Other forms", "EFT / PC / EC stop at Finance — they do not reach you"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "Approve / reject PRs at pending_md",
                "See every request system-wide",
                "Access all reports and analytics",
                "Initiate your own requests",
            ], [
                "## You will NOT typically action:",
                "EFT, Petty Cash, Expense Claims — stop at Finance",
                "GRNs — handled by Finance",
                "User management — Admin's responsibility",
                "## Your focus:",
                "Strategic fit",
                "Value for money",
                "Budget impact on upcoming spend",
                "Supplier and compliance risk",
            ]),
            ("workflow", "Workflow — Your Position", [
                "  [HOD approved] --> pending_finance --> [Finance approved]",
                "                                                 |",
                "                                          pending_md  <-- your queue",
                "                                                 |",
                "                                        approve / reject",
                "                                                 |",
                "                                   +-------------+-----------+",
                "                                   |                         |",
                "                                   v                         v",
                "                               approved                  rejected",
                "                        (Procurement executes)    (initiator notified)",
            ]),
            ("steps", "Task: Review a High-Value Request", [
                "Sidebar → Procurement → Pending Approvals. The list shows only pending_md requests.",
                "Click a request. Read the Approval History — HOD and Finance have commented.",
                "Ask: Is this strategically aligned with company goals?",
                "Ask: Did Finance confirm budget headroom? Are there competing requests this month?",
                "Ask: Is the supplier known and reliable, or new and untested?",
                "Ask: Is there a genuine reason this must be approved now?",
                "If you need more info: call the HOD or Finance directly. Leave the request pending while you get answers.",
            ]),
            ("two_col", "Approve vs Reject", [
                "## To Approve:",
                "Add a comment (even a brief one creates a record)",
                "Click Approve",
                "Status becomes approved",
                "Procurement receives it and proceeds to issue a PO",
            ], [
                "## To Reject:",
                "Write a clear strategic reason",
                "Examples:",
                "-- 'Defer to next quarter — capex freeze in effect.'",
                "-- 'Need competing quotes before I approve.'",
                "-- 'Escalate to Board — exceeds my delegation limit.'",
                "Click Reject",
                "Initiator is notified",
                "## Delegation during absence: tell Admin before you go — they arrange cover.",
            ]),
            ("bullets", "Oversight & Reports", [
                "Sidebar → Reports & Analytics → Reports: summary counts, turnaround times, dept totals",
                "Sidebar → Reports & Analytics → Analytics: dept spend charts, vendor concentration, bottleneck analysis",
                "Export to Excel or PDF for Board packs",
                "## MD weekly habit:",
                "Check pending_md queue — aim to act within 48 hours",
                "Review Analytics once a week for anomalies",
                "Flag unusual vendor patterns to Procurement",
            ]),
        ],
    },

    {
        "file": "05-procurement",
        "label": "Role 5 — Procurement",
        "subtitle": "Vendor management · Quote collection · Adjudication",
        "meta": [("Approval role", "None formal — you execute after MD approval"),
                 ("Key responsibility", "Vendor selection and order execution"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "View every Purchase Requisition across all departments",
                "Add supplier quotes to a requisition",
                "Run an adjudication (three-quote comparison) to select a vendor",
                "Add and maintain the vendor master list",
                "See rejected requisitions to identify recurring patterns",
                "View reports",
            ], [
                "## You CANNOT:",
                "Approve financial requests — that is HOD / Finance / MD",
                "Change amounts or items on an approved PR",
                "Manage users or system settings (Admin only)",
                "## Your responsibilities:",
                "Collect ≥3 quotes for items above the unquoted limit",
                "Run fair adjudications, document the winning vendor and reasoning",
                "Keep the vendor list clean and up to date",
                "Issue purchase orders promptly after MD approval (outside the system)",
            ]),
            ("workflow", "Workflow — Your Position", [
                "  [Initiator] --> pending_hod --> pending_finance --> pending_md --> approved",
                "                                       ^                                |",
                "                                       |                                v",
                "                             YOU collect quotes here         YOU run adjudication,",
                "                             (often while pending_finance)    issue PO outside system",
                "",
                "  Optional sub-flow:",
                "    1. Collect quotes from suppliers (email / PDF).",
                "    2. Upload quotes into the PR detail page.",
                "    3. Run adjudication → select winning vendor → document reasoning.",
                "    4. After MD approves, issue PO to winning vendor.",
            ]),
            ("steps", "Task: Collect & Add Quotes", [
                "Sidebar → Procurement → Approved Requisitions (or watch PRs while still pending_finance).",
                "Open a PR. Scroll to the Quotes section.",
                "Contact at least 3 suppliers and request written quotes (PDF, email, or letterhead).",
                "For each quote: click Add Quote, select the vendor from the master list, enter amount and currency, attach the PDF.",
                "Repeat for each vendor quote.",
                "Once all quotes are in, proceed to adjudication.",
            ]),
            ("steps", "Task: Run Adjudication", [
                "Sidebar → Procurement → Adjudication.",
                "Find the PR and open its adjudication view.",
                "Compare the quotes: price, delivery time, quality/track record.",
                "Select the winning vendor and document the reasoning in the Adjudication Notes.",
                "Submit the adjudication — this records the decision in the audit trail.",
                "Issue the Purchase Order to the winning vendor outside the system.",
            ]),
            ("bullets", "Vendor Management", [
                "Sidebar → Administration → Vendors",
                "## Adding a vendor:",
                "Click Add Vendor → fill name, code, category, contact, tier (1–3), rating (1–5)",
                "## Editing / deactivating:",
                "Click the vendor row → update fields → Save",
                "To deactivate: set Status to inactive — the vendor disappears from dropdowns but historical quotes remain",
                "## Bulk upload:",
                "Download the vendor template → fill columns → click Upload from Excel",
                "## Vendor codes should be stable — do not change a code after it has been used on a quote.",
                "## Quote rules of thumb:",
                "Below the unquoted limit: no quote required, but adding one is always welcome",
                "Above the limit: at least 3 written quotes required",
                "Sole-source: attach a written HOD/Finance justification",
            ]),
        ],
    },

    {
        "file": "06-stores-user",
        "label": "Role 6 — Stores User",
        "subtitle": "Record goods in · Issue goods out · Move stock",
        "meta": [("Account flag", "'Can Access Stores' must be enabled on your profile"),
                 ("Approval role", "None — you create documents, others approve"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "Create Goods Receipt Notes (GRNs) when suppliers deliver",
                "Create Issue Slips when goods leave the warehouse",
                "Create Picking Slips when stock moves between locations (no approval needed)",
                "View live stock levels in the Stock Register",
                "Manage the Stock Items catalogue",
                "Print PDFs of all your documents",
            ], [
                "## You CANNOT:",
                "Approve GRNs or Issue Slips — that is Finance",
                "Create Purchase Requisitions for other departments",
                "Delete a submitted GRN or Issue Slip — ask Finance to reject if wrong",
                "## Your responsibilities:",
                "Record exact quantities — they drive the stock register",
                "Match deliveries to the correct PR",
                "Only issue stock from an approved GRN",
                "Keep the Stock Items list free of duplicates",
            ]),
            ("workflow", "Stores Workflow Diagrams", [
                "GOODS IN (GRN):",
                "  [Supplier delivers] --> [You create GRN] --> pending_approval",
                "                                                      |",
                "                                            Finance approves / rejects",
                "                                                      v",
                "                                                  approved",
                "                                           (stock is now in register)",
                "",
                "GOODS OUT (Issue Slip):",
                "  [You create Issue Slip] --> pending_hod --> pending_finance --> approved",
                "                                                                     |",
                "                                               stock officially deducted from register",
                "",
                "STOCK MOVE (Picking Slip):",
                "  [You create Picking Slip] --> completed immediately (no approval needed)",
            ]),
            ("steps", "Task: Create a Goods Receipt Note (GRN)", [
                "Sidebar → Stores → Create GRN.",
                "Fill the header: GRN Date, Supplier/Vendor, Reference PR (the approved PR this delivery fulfils).",
                "Add each item received: select stock item, enter quantity received, unit, and condition.",
                "If the GRN is for a specific customer, set the Customer field.",
                "Attach the supplier delivery note or invoice.",
                "Click Submit. Status becomes pending_approval. Finance will review and approve.",
                "Only after Finance approval can stock be drawn via Issue Slips.",
            ]),
            ("steps", "Task: Create an Issue Slip", [
                "Confirm the GRN is approved (check Stock Register or GRN list).",
                "Sidebar → Stores → Create Issue Slip.",
                "Fill: Issued To (person / department), Delivery Location, Reference GRN number.",
                "If GRN is customer-reserved, the Customer field must match the GRN.",
                "Add each item: select stock item, enter quantity to issue. The system blocks over-issue automatically.",
                "Click Submit. Goes to pending_hod → pending_finance for approval.",
                "After Finance approval, stock is officially deducted from the register. Hand over the goods.",
            ]),
            ("two_col", "Picking Slips & Stock Register", [
                "## Picking Slip (stock movement)",
                "No approval — completed immediately",
                "Use when physically moving stock between storage areas (e.g. Bay 1 → Bay 3)",
                "Sidebar → Stores → Create Picking Slip",
                "Fill: from location, to location, items and quantities",
                "Submit → status = completed instantly",
                "The register reflects the movement immediately",
            ], [
                "## Real-Time Stock Register",
                "Sidebar → Stores → Real-Time Stock Register",
                "Shows: Stock In (all approved GRNs) minus Stock Out (approved Issue Slips) minus Reserved",
                "Available = what can currently be issued",
                "Use this before creating an Issue Slip to confirm enough stock exists",
                "## Stock Items",
                "Sidebar → Stores → Stock Items",
                "Add or edit catalogue entries (name, unit, description, category)",
                "Avoid duplicate descriptions — one record per item type",
            ]),
        ],
    },

    {
        "file": "07-admin",
        "label": "Role 7 — Administrator",
        "subtitle": "User accounts · Vendors · Clients · Reference data · Emergency interventions",
        "meta": [("Access", "Everything — no hidden menus"),
                 ("Skill level", "Needs basic IT literacy; you are the escalation point for all users"),
                 ("Version", "1.0")],
        "slides": [
            ("two_col", "Who You Are", [
                "## You CAN:",
                "Create, edit, disable, or delete user accounts",
                "Assign roles and departments",
                "Set the 'Can Access Stores' flag",
                "Manage the vendor master list and client list",
                "Set department budgets and FX rates",
                "See every form for support purposes",
                "Reassign stuck approvals",
                "View the audit log",
            ], [
                "## You should NOT:",
                "Approve PRs or financial forms unless you hold a genuine secondary role",
                "Edit historical records to 'fix' them — document via proper channels",
                "Share your admin account — create separate accounts for other admins",
                "## Your core responsibilities:",
                "Keep user master list accurate (role, dept, HOD assignment)",
                "Handle password resets within a few hours",
                "Keep vendor list clean",
                "Preserve the audit log — never manipulate it",
            ]),
            ("steps", "Task: Create a User", [
                "Sidebar → Administration → Users → click Add User.",
                "Fill: Username (e.g. firstname.lastname), Temporary Password (min 8 chars), Full Name, Email, Employee Number.",
                "Assign Role (initiator / hod / procurement / finance / md / admin) and Department.",
                "Tick Is HOD if this person heads their department. Set Supervisor Name to route their requests.",
                "Tick Can Access Stores if the user needs the Stores module.",
                "Click Save. The account is active immediately and flagged to force a password change on first login.",
                "Tell the user their temporary password face-to-face or by phone — the system guides them through the change automatically.",
            ]),
            ("steps", "Task: Reset a User's Password", [
                "Verify the user's identity before acting (face-to-face, by phone, or a known channel — not cold email).",
                "Sidebar → Administration → Users → search and open the user.",
                "Click Reset Password. Enter a new temporary password (min 8 chars). Click Save.",
                "The system hashes the password, marks the account for forced change, and logs the action to the audit trail.",
                "Tell the user the temporary password face-to-face or by phone.",
                "On their next login the system intercepts and forces them to set a new password and 3 security questions.",
            ], "WARNING: A password reset request by cold email is a social-engineering red flag. Always verify identity first."),
            ("two_col", "Vendors & Clients", [
                "## Vendors",
                "Sidebar → Administration → Vendors",
                "Add vendor: name, code, category, contact, tier (1–3), rating (1–5), status",
                "Deactivate: set Status to inactive — historical quotes preserved",
                "Bulk upload: download template → fill → Upload from Excel",
                "Keep vendor codes stable — changing a code after use can break historical references",
            ], [
                "## Clients",
                "Sidebar → Administration → Clients",
                "Clients are external parties whose stock is reserved in your warehouse",
                "Used on GRNs and Issue Slips when the Customer field is set",
                "Add client: name, code, contact, address, status",
                "Deactivate: set Status to inactive",
            ]),
            ("bullets", "Budgets, FX, Audit Log & Emergency Reassignment", [
                "## Budgets",
                "Sidebar → Financial Planning → Budgets — set annual budget per dept in ZMW",
                "Year-end: reset Spent to zero, set new year values; coordinate timing with Finance",
                "## FX Rates",
                "Sidebar → Financial Planning → FX Rates — update only when Finance is unavailable; document the source",
                "## Audit Log",
                "Sidebar → Administration → Audit Log — every approval, rejection, and status change is recorded permanently",
                "You cannot modify or delete audit entries — add a comment on the record if correction context is needed",
                "## Emergency Reassignment (stuck approval)",
                "Identify the stuck request (ask initiator for the form ID)",
                "Decide target: deputy HOD, peer HOD, or escalate to MD",
                "Document the reassignment in the form's Comments with the reason",
                "Notify both the new approver and the initiator",
                "## Admin weekly checklist: password resets, new joiners/leavers, FX check, stuck approvals, audit spot-check",
            ]),
        ],
    },

]

# ── Builder ───────────────────────────────────────────────────────────────────

def build(role):
    prs = new_prs()

    cover_slide(prs,
                role["label"],
                role["subtitle"],
                role["meta"])

    for sdef in role["slides"]:
        kind = sdef[0]
        if kind == "bullets":
            _, title, items = sdef[0], sdef[1], sdef[2]
            note = sdef[3] if len(sdef) > 3 else None
            bullets_slide(prs, title, items, note)

        elif kind == "two_col":
            _, title, left, right = sdef[0], sdef[1], sdef[2], sdef[3]
            # Re-split the flat lists at the heading markers to get titles
            def extract(lst):
                if lst and lst[0].startswith("##"):
                    return lst[0][2:].strip(), lst[1:]
                return "", lst
            lt, li = extract(left)
            rt, ri = extract(right)
            two_col_slide(prs, title, lt, li, rt, ri)

        elif kind == "workflow":
            _, title, lines = sdef[0], sdef[1], sdef[2]
            workflow_slide(prs, title, lines)

        elif kind == "steps":
            _, title, steps = sdef[0], sdef[1], sdef[2]
            intro = sdef[3] if len(sdef) > 3 else None
            steps_slide(prs, title, steps, intro)

    end_slide(prs, role["label"])

    out = os.path.join(os.path.dirname(__file__), f"{role['file']}.pptx")
    prs.save(out)
    print(f"  Saved: {out}")


if __name__ == "__main__":
    print("Generating PPTX presentations...\n")
    for role in ROLES:
        print(f"  {role['label']}")
        build(role)
    print("\nDone.")
